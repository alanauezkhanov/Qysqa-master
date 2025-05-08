from flask import Flask, render_template, request, jsonify, session
import os
import uuid
import json
from werkzeug.utils import secure_filename
import google.generativeai as genai
import PyPDF2
import threading
import time
from datetime import timedelta

# Import libraries for DOCX and PPTX processing
import docx
from pptx import Presentation
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import TextFormatter
import re
import requests
from bs4 import BeautifulSoup
from urllib.parse import urlparse

app = Flask(__name__)
# Use a fixed secret key for production (don't use os.urandom as it changes on restart)
app.secret_key = 'qysqa-fixed-secret-key-for-production'  # Change this to a secure random string in production
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16 MB limit
app.config['PERMANENT_SESSION_LIFETIME'] = timedelta(days=30)  # Set session lifetime to 30 days

# Create upload folder if it doesn't exist
try:
    uploads_dir = app.config['UPLOAD_FOLDER']
    if not os.path.exists(uploads_dir):
        os.makedirs(uploads_dir, mode=0o777, exist_ok=True)
        print(f"Created uploads directory: {uploads_dir}")
    else:
        # Ensure directory has proper permissions
        os.chmod(uploads_dir, 0o777)
        print(f"Uploads directory exists: {uploads_dir}")
except Exception as e:
    print(f"Warning: Could not set up uploads directory: {str(e)}")

# Configure Gemini API
GEMINI_API_KEY = "AIzaSyA9__W2LOpoql4WiEJw0r6a7cCrf2P3P2E"
# AIzaSyBvE04SnLektunSmuCKk0CnzvFSScZupG8
genai.configure(api_key=GEMINI_API_KEY)
model = genai.GenerativeModel('gemini-2.0-flash')

# In-memory storage for user data
users_data = {}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in {
        'pdf', 'txt', 'doc', 'docx', 'ppt', 'pptx'
    }

def process_file(file_path, file_name, user_id):
    """Extract text from file and save to user's sources"""
    
    text_content = ""
    try:
        if file_name.lower().endswith('.txt'):
            print(f"Processing TXT file: {file_path}")
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text_content = f.read()
                print(f"Successfully read text file with length: {len(text_content)}")
            except Exception as e:
                print(f"Error reading text file: {str(e)}")
                # Try alternative encoding
                try:
                    with open(file_path, 'r', encoding='cp1251', errors='ignore') as f:
                        text_content = f.read()
                    print(f"Successfully read text file with cp1251 encoding, length: {len(text_content)}")
                except Exception as e2:
                    print(f"Error reading text file with alternative encoding: {str(e2)}")
                    # Last resort - read as binary and decode
                    with open(file_path, 'rb') as f:
                        binary_content = f.read()
                        text_content = binary_content.decode('utf-8', errors='ignore')
                    print(f"Read file as binary, decoded length: {len(text_content)}")
        elif file_name.lower().endswith('.pdf'):
            print(f"Processing PDF file: {file_path}")
            try:
                # Extract text from PDF
                text_content = extract_text_from_pdf(file_path)
                print(f"Successfully extracted text from PDF, length: {len(text_content)}")
            except Exception as e:
                print(f"Error extracting text from PDF: {str(e)}")
                text_content = f"Error extracting content from {file_name}: {str(e)}"
        elif file_name.lower().endswith('.docx'):
            print(f"Processing DOCX file: {file_path}")
            try:
                # Extract text from DOCX
                text_content = extract_text_from_docx(file_path)
                print(f"Successfully extracted text from DOCX, length: {len(text_content)}")
            except Exception as e:
                print(f"Error extracting text from DOCX: {str(e)}")
                text_content = f"Error extracting content from {file_name}: {str(e)}"
        elif file_name.lower().endswith('.pptx'):
            print(f"Processing PPTX file: {file_path}")
            try:
                # Extract text from PPTX
                text_content = extract_text_from_pptx(file_path)
                print(f"Successfully extracted text from PPTX, length: {len(text_content)}")
            except Exception as e:
                print(f"Error extracting text from PPTX: {str(e)}")
                text_content = f"Error extracting content from {file_name}: {str(e)}"
        else:
            text_content = f"Content placeholder for {file_name} - This file type is not yet supported for content extraction."
        
        # Add to user's sources
        if user_id not in users_data:
            users_data[user_id] = {
                "sources": [],
                "chat_history": []
            }
        
        source_id = str(uuid.uuid4())
        users_data[user_id]["sources"].append({
            "id": source_id,
            "name": file_name,
            "content": text_content
        })
        
        return source_id
    except Exception as e:
        print(f"Error in process_file: {str(e)}")
        raise Exception(f"Could not process file: {str(e)}")

def extract_text_from_pdf(file_path):
    """Extract text from a PDF file"""
    text = ""
    try:
        # Open the PDF file
        with open(file_path, 'rb') as file:
            # Create a PDF reader object
            pdf_reader = PyPDF2.PdfReader(file)
            
            # Get the number of pages
            num_pages = len(pdf_reader.pages)
            print(f"PDF has {num_pages} pages")
            
            # Extract text from each page
            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n\n"
                else:
                    print(f"Warning: No text extracted from page {page_num+1}")
                    
            if not text.strip():
                print("Warning: No text extracted from PDF")
                text = "This PDF does not contain extractable text. It may be scanned or contain only images."
                
    except Exception as e:
        print(f"Error in extract_text_from_pdf: {str(e)}")
        raise Exception(f"Failed to extract text from PDF: {str(e)}")
        
    return text

def extract_text_from_docx(file_path):
    """Extract text from a DOCX file"""
    text = ""
    try:
        # Load the DOCX document
        doc = docx.Document(file_path)
        
        # Extract text from paragraphs
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        
        # Extract text from tables
        tables_text = []
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                if row_text:
                    tables_text.append(row_text)
        
        # Combine paragraph and table text
        if paragraphs:
            text += "Document Text:\n" + "\n".join(paragraphs) + "\n\n"
        
        if tables_text:
            text += "Tables:\n" + "\n".join(tables_text)
            
        if not text.strip():
            print("Warning: No text extracted from DOCX")
            text = "This DOCX file does not contain extractable text. It may contain only images or other non-text elements."
            
    except Exception as e:
        print(f"Error in extract_text_from_docx: {str(e)}")
        raise Exception(f"Failed to extract text from DOCX: {str(e)}")
        
    return text

def extract_text_from_pptx(file_path):
    """Extract text from a PPTX file"""
    text = ""
    try:
        # Load the PPTX presentation
        pres = Presentation(file_path)
        
        # Process each slide
        for slide_num, slide in enumerate(pres.slides, 1):
            slide_text = []
            
            # Get slide title if exists
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_text.append(f"Title: {slide.shapes.title.text.strip()}")
            
            # Extract text from each shape in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip() and (not slide.shapes.title or shape != slide.shapes.title):
                    slide_text.append(shape.text.strip())
            
            # Add slide text to overall text
            if slide_text:
                text += f"Slide {slide_num}:\n" + "\n".join(slide_text) + "\n\n"
        
        if not text.strip():
            print("Warning: No text extracted from PPTX")
            text = "This PPTX file does not contain extractable text. It may contain only images or other non-text elements."
            
    except Exception as e:
        print(f"Error in extract_text_from_pptx: {str(e)}")
        raise Exception(f"Failed to extract text from PPTX: {str(e)}")
        
    return text

def generate_summary(user_id):
    """Generate summary from all user sources"""
    if user_id not in users_data or not users_data[user_id]["sources"]:
        return "No content available to summarize."
    
    # Combine all source content
    all_content = "\n\n".join([src["content"] for src in users_data[user_id]["sources"]])
    
    # Generate summary using Gemini
    prompt = f"""
    Please analyze the following educational content and provide:
    1. A concise summary of the main topics
    2. Key concepts and terminology highlighted
    3. A brief overview of how these concepts relate to each other
    
    FORMAT YOUR RESPONSE USING MARKDOWN:
    - Use ## or ### for section headings
    - Use **bold** for key terms and important concepts
    - Use *italic* for emphasis
    - Use bullet points or numbered lists for structured information
    - Use tables if presenting structured data
    - Use > blockquotes for highlighting important information or quotes
    - Organize the summary with clear sections and hierarchy
    
    Content:
    {all_content[:4000]}  # Limit content length for API
    """
    
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        print(f"Error generating summary: {e}")
        return "Could not generate summary. Please try again later."

@app.route('/')
def index():
    # Create a unique user ID if not exists
    if 'user_id' not in session:
        session['user_id'] = str(uuid.uuid4())
        # Make session permanent
        session.permanent = True
    
    user_id = session['user_id']
    if user_id not in users_data:
        users_data[user_id] = {
            "sources": [],
            "chat_history": []
        }
    
    return render_template('index.html')

@app.route('/api/upload', methods=['POST'])
def upload_file():
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No file part"}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "No selected file"}), 400
        
        print(f"Received file: {file.filename}, mimetype: {file.mimetype}")
        
        if file and allowed_file(file.filename):
            user_id = session.get('user_id')
            if not user_id:
                return jsonify({"error": "User session expired"}), 401
            
            try:
                filename = secure_filename(file.filename)
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{user_id}_{filename}")
                file.save(file_path)
                
                print(f"File saved to {file_path}")
                
                # Process the file
                try:
                    source_id = process_file(file_path, filename, user_id)
                    
                    return jsonify({
                        "success": True,
                        "message": "File uploaded and processed successfully",
                        "filename": filename,
                        "source_id": source_id
                    })
                except Exception as e:
                    error_msg = f"Error processing file: {str(e)}"
                    print(error_msg)
                    return jsonify({"error": error_msg}), 500
            except Exception as e:
                error_msg = f"Error saving file: {str(e)}"
                print(error_msg)
                return jsonify({"error": error_msg}), 500
        
        return jsonify({"error": f"File type not allowed: {file.filename}"}), 400
    except Exception as e:
        # Catch-all exception handler to ensure we always return valid JSON
        error_msg = f"Unexpected error during file upload: {str(e)}"
        print(error_msg)
        return jsonify({"error": error_msg}), 500

@app.route('/api/sources', methods=['GET'])
def get_sources():
    user_id = session.get('user_id')
    if not user_id or user_id not in users_data:
        return jsonify([])
    
    sources = []
    for source in users_data[user_id]["sources"]:
        sources.append({
            "id": source["id"],
            "name": source["name"]
        })
    
    return jsonify(sources)

@app.route('/api/summarize', methods=['GET'])
def summarize():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "User session expired"}), 401
    
    summary = generate_summary(user_id)
    return jsonify({"summary": summary})

@app.route('/api/chat', methods=['POST'])
def chat():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "User session expired"}), 401
    
    data = request.json
    if not data or 'message' not in data:
        return jsonify({"error": "No message provided"}), 400
    
    user_message = data['message']
    
    # Get the selected sources
    selected_sources = data.get('sources', [])
    print(f"Selected sources: {selected_sources}")
    
    # Get content from selected sources
    context = ""
    if selected_sources:
        # Only use selected sources if they're specified
        for source_id in selected_sources:
            found = False
            for source in users_data[user_id]["sources"]:
                if source["id"] == source_id:
                    context += f"\nContent from {source['name']}:\n{source['content']}\n"
                    found = True
            if not found:
                print(f"Warning: Source ID {source_id} not found")
        
        if not context.strip():
            return jsonify({
                "message": "Я не могу ответить на ваш вопрос, так как вы не выбрали ни одного источника или выбранные источники недоступны. Пожалуйста, выберите один или несколько источников."
            })
    else:
        # If no sources selected, inform the user instead of using all sources
        return jsonify({
            "message": "Пожалуйста, выберите один или несколько источников, чтобы я мог ответить на ваш вопрос. Для этого отметьте галочками нужные источники в левой панели."
        })
    
    # Prepare prompt with context, instructing to use Markdown
    prompt = f"""
    You are an educational assistant helping students understand their course materials.
    Use the following context to answer the question:
    
    {context[:4000]}  # Limit context length for API
    
    Student question: {user_message}
    
    Provide a concise, helpful answer based on the provided materials.
    
    FORMAT YOUR RESPONSE USING MARKDOWN:
    - Use **bold** or *italic* for emphasis
    - Use headings (## or ###) for section titles
    - Use bullet points or numbered lists where appropriate
    - Use code blocks for any code or formatted text
    - Use tables if presenting tabular data
    - Use blockquotes for highlighting important information
    
    Be sure to make your response visually structured and easy to read.
    """
    
    try:
        response = model.generate_content(prompt)
        ai_message = response.text
        
        # Record in chat history
        if user_id not in users_data:
            users_data[user_id] = {"sources": [], "chat_history": []}
        
        users_data[user_id]["chat_history"].append({
            "user": user_message,
            "assistant": ai_message
        })
        
        return jsonify({
            "message": ai_message
        })
        
    except Exception as e:
        print(f"Error in chat generation: {e}")
        return jsonify({
            "error": "Failed to generate response",
            "message": "I'm having trouble connecting right now. Please try again later."
        }), 500

@app.route('/api/generate-test', methods=['POST'])
def generate_test():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "User session expired"}), 401
    
    data = request.json
    num_closed_questions = int(data.get('num_closed_questions', 5))  # Convert to integer
    num_open_questions = int(data.get('num_open_questions', 0))  # Convert to integer
    difficulty = data.get('difficulty', 'medium')
    
    total_questions = num_closed_questions + num_open_questions
    print(f"Generating test with {num_closed_questions} closed questions and {num_open_questions} open questions at {difficulty} difficulty")
    
    # Get selected sources
    selected_sources = data.get('sources', [])
    print(f"Selected sources for test generation: {selected_sources}")
    
    # Get sources from user data
    if not selected_sources:
        # Use all sources if none selected
        print("No sources selected, using all available sources")
        sources_to_use = users_data[user_id]["sources"]
    else:
        # Use only selected sources
        sources_to_use = []
        for source_id in selected_sources:
            for source in users_data[user_id]["sources"]:
                if source["id"] == source_id:
                    sources_to_use.append(source)
                    print(f"Added source '{source['name']}' to test generation list")
                    break
    
    if not sources_to_use:
        return jsonify({"error": "No valid sources available to generate test"}), 400
    
    print(f"Using {len(sources_to_use)} sources for test generation")
    
    # Calculate questions per source (distribute evenly with at least 1 per source if possible)
    total_sources = len(sources_to_use)
    
    # Initialize questions list
    all_questions = []
    
    # First handle closed questions if any
    if num_closed_questions > 0:
        closed_questions = generate_closed_questions(sources_to_use, num_closed_questions, difficulty)
        all_questions.extend(closed_questions)
    
    # Then handle open questions if any
    if num_open_questions > 0:
        open_questions = generate_open_questions(sources_to_use, num_open_questions, difficulty)
        all_questions.extend(open_questions)
    
    # If we didn't get any questions, return an error
    if not all_questions:
        return jsonify({
            "error": "Failed to generate any valid questions",
            "message": "The AI couldn't create test questions from your materials. Try different sources or simplify your content."
        }), 500
    
    print(f"Total questions generated: {len(all_questions)}")
    
    return jsonify({
        "success": True,
        "questions": all_questions
    })

def generate_closed_questions(sources_to_use, num_questions, difficulty):
    """Generate closed (multiple-choice) questions"""
    total_sources = len(sources_to_use)
    
    if total_sources > num_questions:
        # If we have more sources than questions, we'll use a subset of sources
        sources_to_use = sources_to_use[:num_questions]
        questions_per_source = [1] * len(sources_to_use)
        total_sources = len(sources_to_use)
    else:
        base_questions = num_questions // total_sources
        extra_questions = num_questions % total_sources
        questions_per_source = [base_questions] * total_sources
        # Distribute extra questions
        for i in range(extra_questions):
            questions_per_source[i] += 1
    
    print(f"Closed question distribution: {questions_per_source} across {total_sources} sources")
    
    # Build a prompt that explicitly requests specific questions from each source
    main_prompt = f"""
    You are creating a multiple-choice test with {num_questions} closed questions at {difficulty} difficulty level.
    
    Each question must have:
    1. A clear, concise question
    2. Four multiple choice options labeled A, B, C, D
    3. The correct answer letter
    4. A brief explanation of why it's correct (use Markdown formatting in the explanation)
    5. Source information
    
    FORMAT THE EXPLANATION USING MARKDOWN:
    - Use **bold** for emphasis on key points
    - Use *italic* for terms or concepts
    - Use bullet points for multiple points
    - Use > blockquote for important notes
    
    Format each question as follows:
    
    Question X: [The question text]
    A. [Option A]
    B. [Option B]
    C. [Option C]
    D. [Option D]
    Correct Answer: [Letter]
    Explanation: [Why this answer is correct - use Markdown formatting]
    Source: [Name of the source]
    """
    
    # Generate questions source by source
    all_questions = []
    
    try:
        for i, source in enumerate(sources_to_use):
            source_content = source['content']
            source_name = source['name']
            questions_needed = questions_per_source[i]
            
            print(f"Generating {questions_needed} closed questions from source: {source_name}")
            
            source_prompt = f"""
{main_prompt}

I need you to generate {questions_needed} questions specifically from this content:

SOURCE NAME: {source_name}
CONTENT:
{source_content[:4000]}  # Limit content length for API

Make exactly {questions_needed} questions from this material.
Each question MUST be based on the content above and include "Source: {source_name}" at the end.
"""
            
            # Generate questions for this source
            try:
                print(f"Sending prompt to Gemini API for source: {source_name}")
                response = model.generate_content(source_prompt)
                source_response = response.text
                print(f"Response received for source {source_name}, length: {len(source_response)}")
                
                # Parse questions from this response
                source_questions = parse_questions(source_response, default_source=source_name)
                print(f"Parsed {len(source_questions)} questions from source: {source_name}")
                
                # Mark questions as closed type
                for q in source_questions:
                    q['type'] = 'closed'
                
                # Add to our collection
                all_questions.extend(source_questions)
                
            except Exception as e:
                print(f"Error generating closed questions for source {source_name}: {e}")
                continue
        
        # If we didn't get any questions, try one more time with all sources combined
        if not all_questions:
            print("No questions generated from individual sources, trying with combined sources")
            
            # Build combined content
            combined_content = ""
            for source in sources_to_use:
                combined_content += f"\n--- {source['name']} ---\n{source['content'][:1000]}\n\n"
            
            fallback_prompt = f"""
{main_prompt}

I need you to generate {num_questions} questions from the following sources:

{combined_content}

Make sure to indicate which source each question comes from using "Source: [source name]" at the end of each question.
"""
            
            response = model.generate_content(fallback_prompt)
            questions = parse_questions(response.text)
            
            # Mark questions as closed type
            for q in questions:
                q['type'] = 'closed'
                
            all_questions.extend(questions)
            
    except Exception as e:
        print(f"Error in generate_closed_questions: {e}")
        import traceback
        traceback.print_exc()
        
    return all_questions

def generate_open_questions(sources_to_use, num_questions, difficulty):
    """Generate open-ended questions"""
    total_sources = len(sources_to_use)
    
    if total_sources > num_questions:
        # If we have more sources than questions, we'll use a subset of sources
        sources_to_use = sources_to_use[:num_questions]
        questions_per_source = [1] * len(sources_to_use)
        total_sources = len(sources_to_use)
    else:
        base_questions = num_questions // total_sources
        extra_questions = num_questions % total_sources
        questions_per_source = [base_questions] * total_sources
        # Distribute extra questions
        for i in range(extra_questions):
            questions_per_source[i] += 1
    
    print(f"Open question distribution: {questions_per_source} across {total_sources} sources")
    
    # Build a prompt that explicitly requests specific questions from each source
    main_prompt = f"""
    You are creating an exam with {num_questions} open-ended questions at {difficulty} difficulty level.
    
    Each question must have:
    1. A clear, thought-provoking question that requires a detailed answer
    2. A model answer that would receive a perfect score (10/10)
    3. Evaluation criteria for scoring student answers from 1-10
    4. Source information
    
    FORMAT THE MODEL ANSWER AND CRITERIA USING MARKDOWN:
    - Use **bold** for emphasis on key points
    - Use *italic* for terms or concepts
    - Use bullet points for multiple points
    
    Format each question as follows:
    
    Question: [The question text]
    Model Answer: [A detailed perfect answer that would score 10/10]
    Evaluation Criteria: [How to score responses from 1-10]
    Source: [Name of the source]
    """
    
    # Generate questions source by source
    all_questions = []
    
    try:
        for i, source in enumerate(sources_to_use):
            source_content = source['content']
            source_name = source['name']
            questions_needed = questions_per_source[i]
            
            print(f"Generating {questions_needed} open questions from source: {source_name}")
            
            source_prompt = f"""
{main_prompt}

I need you to generate {questions_needed} open-ended questions specifically from this content:

SOURCE NAME: {source_name}
CONTENT:
{source_content[:4000]}  # Limit content length for API

Make exactly {questions_needed} questions from this material.
Each question MUST be based on the content above and include "Source: {source_name}" at the end.

Each evaluation criteria should clearly explain what makes a response score:
- 1-3 (poor understanding, major misconceptions)
- 4-6 (basic understanding but incomplete)
- 7-8 (good understanding with minor omissions)
- 9-10 (excellent understanding, comprehensive answer)
"""
            
            # Generate questions for this source
            try:
                print(f"Sending prompt to Gemini API for source: {source_name}")
                response = model.generate_content(source_prompt)
                source_response = response.text
                print(f"Response received for source {source_name}, length: {len(source_response)}")
                
                # Parse open questions from this response
                source_questions = parse_open_questions(source_response, default_source=source_name)
                print(f"Parsed {len(source_questions)} open questions from source: {source_name}")
                
                # Add to our collection
                all_questions.extend(source_questions)
                
            except Exception as e:
                print(f"Error generating open questions for source {source_name}: {e}")
                continue
        
        # If we didn't get any questions, try one more time with all sources combined
        if not all_questions:
            print("No open questions generated from individual sources, trying with combined sources")
            
            # Build combined content
            combined_content = ""
            for source in sources_to_use:
                combined_content += f"\n--- {source['name']} ---\n{source['content'][:1000]}\n\n"
            
            fallback_prompt = f"""
{main_prompt}

I need you to generate {num_questions} open-ended questions from the following sources:

{combined_content}

Make sure to indicate which source each question comes from using "Source: [source name]" at the end of each question.

IMPORTANT: Format each question following this EXACT structure:
Question: [question text]
Model Answer: [model answer text]
Evaluation Criteria: [evaluation criteria text]
Source: [source name]

This exact formatting is required for the questions to be parsed correctly.
"""
            
            try:
                print("Sending fallback prompt to Gemini API for combined sources")
                response = model.generate_content(fallback_prompt)
                combined_response = response.text
                print(f"Response received for combined sources, length: {len(combined_response)}")
                
                # Parse open questions from the response
                all_questions = parse_open_questions(combined_response)
                print(f"Parsed {len(all_questions)} open questions from combined sources")
                
                # If still no questions, create simple placeholder questions
                if not all_questions and num_questions > 0:
                    print("Creating placeholder open questions")
                    for i in range(num_questions):
                        source = sources_to_use[i % len(sources_to_use)]
                        source_name = source['name']
                        
                        # Extract a sentence from the source to use as a question
                        content_sample = source['content'][:2000]
                        sentences = re.split(r'[.!?]', content_sample)
                        question_base = sentences[min(i, len(sentences)-1)] if sentences else "Explain this topic"
                        
                        all_questions.append({
                            'type': 'open',
                            'question': f"Please explain the following in detail: {question_base.strip()}?",
                            'model_answer': f"A complete answer would include key information from {source_name}.",
                            'evaluation_criteria': "Score 1-3: Minimal understanding. Score 4-6: Basic understanding. Score 7-8: Good understanding. Score 9-10: Excellent, comprehensive answer.",
                            'source': source_name
                        })
                    print(f"Created {len(all_questions)} placeholder open questions")
            except Exception as e:
                print(f"Error in fallback open question generation: {e}")
                import traceback
                traceback.print_exc()
        
    except Exception as e:
        print(f"Error in generate_open_questions: {e}")
        import traceback
        traceback.print_exc()
        
    return all_questions

@app.route('/api/diagnostic', methods=['GET'])
def diagnostic():
    """Diagnostic endpoint to check server functionality"""
    try:
        # Check uploads directory
        uploads_dir = app.config['UPLOAD_FOLDER']
        uploads_exists = os.path.exists(uploads_dir)
        uploads_writable = os.access(uploads_dir, os.W_OK) if uploads_exists else False
        
        # Check session
        session_working = 'user_id' in session
        
        # Try to write a test file
        test_file_path = os.path.join(uploads_dir, 'test_file.txt')
        test_file_success = False
        try:
            with open(test_file_path, 'w') as f:
                f.write('Test file for diagnostics')
            test_file_success = True
            os.remove(test_file_path)  # Clean up
        except Exception as e:
            test_file_error = str(e)
        
        return jsonify({
            "status": "ok",
            "uploads_directory": {
                "path": uploads_dir,
                "exists": uploads_exists,
                "writable": uploads_writable
            },
            "session_working": session_working,
            "test_file_write": {
                "success": test_file_success,
                "error": test_file_error if not test_file_success else None
            },
            "memory_usage": {
                "sources": len(users_data),
                "sources_details": {uid: len(data["sources"]) for uid, data in users_data.items()}
            }
        })
    except Exception as e:
        return jsonify({
            "status": "error",
            "error": str(e)
        })

def parse_questions(text, default_source="Unknown"):
    """Parse questions from AI response text"""
    questions = []
    current_question = {}
    
    # Split by lines and process
    lines = text.split('\n')
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        if line.startswith('Question') or (len(questions) > 0 and not current_question and line):
            # Save previous question if it exists
            if current_question and 'question' in current_question and 'options' in current_question and len(current_question['options']) > 0:
                questions.append(current_question)
            
            # Start new question
            current_question = {
                'question': line.split(':', 1)[1].strip() if ':' in line else line,
                'options': [],
                'correct_answer': '',
                'explanation': '',
                'source': default_source
            }
        
        # Parse options
        elif line.startswith(('A.', 'B.', 'C.', 'D.')) and current_question:
            option_letter = line[0]
            option_text = line[2:].strip()
            current_question['options'].append(option_text)
        
        # Parse correct answer
        elif line.startswith('Correct Answer:') and current_question:
            current_question['correct_answer'] = line.split(':', 1)[1].strip()
        
        # Parse explanation
        elif line.startswith('Explanation:') and current_question:
            current_question['explanation'] = line.split(':', 1)[1].strip()
            
        # Parse source
        elif line.startswith('Source:') and current_question:
            current_question['source'] = line.split(':', 1)[1].strip()
        
        # Append to explanation if it's a continuation and not another field
        elif current_question and 'explanation' in current_question and current_question['explanation'] and not line.startswith(('Question', 'A.', 'B.', 'C.', 'D.', 'Correct Answer:', 'Source:')):
            current_question['explanation'] += ' ' + line
    
    # Add the last question if it exists
    if current_question and 'question' in current_question and 'options' in current_question and len(current_question['options']) > 0:
        questions.append(current_question)
    
    return questions

def parse_open_questions(text, default_source="Unknown"):
    """Parse open-ended questions from AI response text"""
    questions = []
    current_question = {}
    
    # Split by lines and process
    lines = text.split('\n')
    current_section = None
    
    # Debug flag to log parsing steps
    debug = False
    
    if debug:
        print("Parsing open questions, total lines:", len(lines))
    
    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue
        
        # Debug output
        if debug and i < 10:
            print(f"Line {i}: {line[:50]} ... (section: {current_section})")
        
        # Check for new question start - handle various formats
        if (re.match(r'^question\s*(\d+)?\s*:.*', line.lower()) or 
            re.match(r'^question\s*(\d+)?$', line.lower()) or
            re.match(r'^\d+\.\s*question:.*', line.lower()) or
            (line.startswith('Question') and not current_question)):
            
            # Save previous question if it exists and is complete
            if current_question and 'question' in current_question and 'model_answer' in current_question and 'evaluation_criteria' in current_question:
                if debug:
                    print(f"Adding complete question: {current_question.get('question', '')[:30]}...")
                questions.append(current_question)
            elif current_question:
                if debug:
                    print(f"Skipping incomplete question: {current_question.get('question', '')[:30]}...")
            
            # Start new question
            question_text = ""
            if ':' in line:
                question_text = line.split(':', 1)[1].strip()
            else:
                # If this line is just a header, check the next line for the actual question
                for next_line in lines[i+1:]:
                    next_line = next_line.strip()
                    if next_line and not next_line.lower().startswith(('model answer:', 'evaluation criteria:')):
                        question_text = next_line
                        break
            
            current_question = {
                'question': question_text,
                'model_answer': '',
                'evaluation_criteria': '',
                'source': default_source,
                'type': 'open'  # Mark as open-ended question
            }
            current_section = 'question'
            
            if debug:
                print(f"Started new question: {question_text[:30]}...")
            
        # Check for model answer section with various formats
        elif (line.lower().startswith('model answer:') or 
              line.lower() == 'model answer' or 
              line.lower().startswith('answer:') or
              re.match(r'^model\s+answer\s*:', line.lower())):
            
            current_section = 'model_answer'
            if ':' in line:
                current_question['model_answer'] = line.split(':', 1)[1].strip()
            
            if debug:
                print(f"Started model answer section: {current_question.get('model_answer', '')[:30]}...")
                
        # Check for evaluation criteria section with various formats
        elif (line.lower().startswith('evaluation criteria:') or 
              line.lower() == 'evaluation criteria' or
              line.lower().startswith('criteria:') or
              re.match(r'^evaluation\s+criteria\s*:', line.lower())):
            
            current_section = 'evaluation_criteria'
            if ':' in line:
                current_question['evaluation_criteria'] = line.split(':', 1)[1].strip()
            
            if debug:
                print(f"Started evaluation criteria section: {current_question.get('evaluation_criteria', '')[:30]}...")
            
        # Check for source information
        elif line.lower().startswith('source:'):
            current_question['source'] = line.split(':', 1)[1].strip()
            current_section = 'source'
            
            if debug:
                print(f"Added source: {current_question.get('source', '')[:30]}...")
            
        # Append content to current section if it's not a new section header
        elif current_section and current_question:
            # Check if this line is a start of a new section
            is_new_section = any([
                re.match(r'^question\s*(\d+)?\s*:.*', line.lower()),
                line.lower().startswith('model answer:'),
                line.lower() == 'model answer',
                line.lower().startswith('answer:'),
                line.lower().startswith('evaluation criteria:'),
                line.lower() == 'evaluation criteria',
                line.lower().startswith('criteria:'),
                line.lower().startswith('source:')
            ])
            
            if not is_new_section:
                if current_section in current_question:
                    if current_question[current_section]:  # If there's already content, add a space before appending
                        current_question[current_section] += ' ' + line
                    else:
                        current_question[current_section] = line
    
    # Add the last question if it exists and is complete
    if current_question and 'question' in current_question and 'model_answer' in current_question and 'evaluation_criteria' in current_question:
        if debug:
            print(f"Adding final question: {current_question.get('question', '')[:30]}...")
        questions.append(current_question)
    
    if debug:
        print(f"Total parsed open questions: {len(questions)}")
    
    return questions

def extract_youtube_id(url):
    """Extract YouTube video ID from URL"""
    youtube_regex = (
        r'(https?://)?(www\.)?'
        '(youtube|youtu|youtube-nocookie)\.(com|be)/'
        '(watch\?v=|embed/|v/|.+\?v=)?([^&=%\?]{11})')
    match = re.match(youtube_regex, url)
    if match:
        return match.group(6)
    return None

def extract_transcript_from_youtube(video_id):
    """Extract transcript from YouTube video"""
    try:
        transcript_list = YouTubeTranscriptApi.get_transcript(video_id, languages=['en', 'ru', 'kk'])
        
        # Instead of using the formatter, manually combine transcript parts
        transcript_text = ""
        for entry in transcript_list:
            # Each entry contains 'text', 'start', and 'duration'
            if 'text' in entry:
                transcript_text += entry['text'] + " "
        
        if not transcript_text.strip():
            raise Exception("No transcript text was extracted")
            
        return transcript_text
    except Exception as e:
        print(f"Error extracting YouTube transcript: {str(e)}")
        raise Exception(f"Could not extract transcript: {str(e)}")

def get_youtube_video_title(video_id):
    """Try to get the title of a YouTube video"""
    try:
        # Use a simple HTML fetch and regex to extract the title
        # (Using YouTube Data API would be better but requires API key setup)
        response = requests.get(f"https://www.youtube.com/watch?v={video_id}")
        if response.status_code == 200:
            html_content = response.text
            title_match = re.search(r'<title>(.*?) - YouTube</title>', html_content)
            if title_match:
                return title_match.group(1)
        return f"YouTube Video ({video_id})"
    except Exception as e:
        print(f"Error fetching video title: {e}")
        return f"YouTube Video ({video_id})"

@app.route('/api/add-youtube', methods=['POST'])
def add_youtube():
    """Add a YouTube video as a source by extracting its transcript"""
    try:
        user_id = session.get('user_id')
        if not user_id:
            return jsonify({"error": "User session expired"}), 401
        
        data = request.json
        if not data or 'youtube_url' not in data:
            return jsonify({"error": "No YouTube URL provided"}), 400
        
        youtube_url = data['youtube_url']
        video_id = extract_youtube_id(youtube_url)
        
        if not video_id:
            return jsonify({"error": "Invalid YouTube URL"}), 400
        
        print(f"Processing YouTube video: {youtube_url} (ID: {video_id})")
        
        try:
            # Extract transcript
            transcript = extract_transcript_from_youtube(video_id)
            
            # Try to get video title
            video_title = get_youtube_video_title(video_id)
            
            # Add to user's sources
            if user_id not in users_data:
                users_data[user_id] = {
                    "sources": [],
                    "chat_history": []
                }
            
            source_id = str(uuid.uuid4())
            users_data[user_id]["sources"].append({
                "id": source_id,
                "name": video_title,
                "content": transcript,
                "type": "youtube",
                "video_id": video_id
            })
            
            return jsonify({
                "success": True,
                "message": "YouTube transcript extracted and added successfully",
                "video_id": video_id,
                "source_id": source_id
            })
            
        except Exception as e:
            error_msg = f"Error processing YouTube video: {str(e)}"
            print(error_msg)
            return jsonify({"error": error_msg}), 500
        
    except Exception as e:
        error_msg = f"Unexpected error processing YouTube URL: {str(e)}"
        print(error_msg)
        return jsonify({"error": error_msg}), 500

@app.route('/api/add-clipboard', methods=['POST'])
def add_clipboard():
    """Add text from clipboard as a source"""
    try:
        user_id = session.get('user_id')
        if not user_id:
            return jsonify({"error": "User session expired"}), 401
        
        data = request.json
        if not data or 'text' not in data or not data['text'].strip():
            return jsonify({"error": "No text provided"}), 400
        
        text_content = data['text'].strip()
        source_name = data.get('name', 'Clipboard Text')
        
        print(f"Adding clipboard text as source: {source_name} ({len(text_content)} characters)")
        
        # Add to user's sources
        if user_id not in users_data:
            users_data[user_id] = {
                "sources": [],
                "chat_history": []
            }
        
        source_id = str(uuid.uuid4())
        users_data[user_id]["sources"].append({
            "id": source_id,
            "name": source_name,
            "content": text_content,
            "type": "text"
        })
        
        return jsonify({
            "success": True,
            "message": "Text added successfully",
            "source_id": source_id
        })
        
    except Exception as e:
        error_msg = f"Unexpected error adding clipboard text: {str(e)}"
        print(error_msg)
        return jsonify({"error": error_msg}), 500

def extract_text_from_website(url):
    """Extract main content text from a website"""
    try:
        # Add headers to mimic a browser request
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        
        # Fetch the webpage
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()  # Raise an exception for 4XX/5XX responses
        
        # Parse the HTML content
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # Remove script and style elements
        for script_or_style in soup(['script', 'style', 'header', 'footer', 'nav', 'aside']):
            script_or_style.decompose()
        
        # Extract site title for source name
        title = soup.title.string if soup.title else url
        
        # Try to extract the main content by focusing on likely content containers
        main_content = soup.find('main') or soup.find('article') or soup.find('div', class_=lambda c: c and any(x in str(c).lower() for x in ['content', 'main', 'article', 'post']))
        
        if main_content:
            # If a main content container was found, extract text from it
            text = main_content.get_text(separator='\n', strip=True)
        else:
            # If no main content container was found, extract text from the body
            text = soup.get_text(separator='\n', strip=True)
        
        # Clean up the text
        text = re.sub(r'\n+', '\n', text)  # Replace multiple newlines with single newline
        text = re.sub(r'\s+', ' ', text)    # Replace multiple spaces with single space
        
        # If the text is too long, truncate it
        if len(text) > 100000:
            text = text[:100000] + '\n\n[Content truncated due to length...]'
        
        return {
            'text': text,
            'title': title
        }
    except Exception as e:
        print(f"Error extracting text from website: {str(e)}")
        raise Exception(f"Failed to extract content from website: {str(e)}")

@app.route('/api/add-website', methods=['POST'])
def add_website():
    """Add content from a website as a source"""
    try:
        user_id = session.get('user_id')
        if not user_id:
            return jsonify({"error": "User session expired"}), 401
        
        data = request.json
        if not data or 'url' not in data:
            return jsonify({"error": "No URL provided"}), 400
        
        url = data['url'].strip()
        
        # Basic URL validation
        if not url.startswith(('http://', 'https://')):
            url = 'https://' + url
        
        try:
            # Parse the URL to extract domain for source name
            parsed_url = urlparse(url)
            domain = parsed_url.netloc
            
            print(f"Processing website: {url}")
            
            # Extract text from website
            result = extract_text_from_website(url)
            
            text_content = result['text']
            title = result['title']
            
            # Create source name from title and domain
            source_name = f"{title} ({domain})"
            
            # Add to user's sources
            if user_id not in users_data:
                users_data[user_id] = {
                    "sources": [],
                    "chat_history": []
                }
            
            source_id = str(uuid.uuid4())
            users_data[user_id]["sources"].append({
                "id": source_id,
                "name": source_name,
                "content": text_content,
                "type": "website",
                "url": url
            })
            
            return jsonify({
                "success": True,
                "message": "Website content extracted and added successfully",
                "source_id": source_id,
                "url": url
            })
            
        except Exception as e:
            error_msg = f"Error processing website: {str(e)}"
            print(error_msg)
            return jsonify({"error": error_msg}), 500
        
    except Exception as e:
        error_msg = f"Unexpected error processing website URL: {str(e)}"
        print(error_msg)
        return jsonify({"error": error_msg}), 500

@app.route('/api/evaluate-answer', methods=['POST'])
def evaluate_answer():
    """Evaluate an open-ended question answer"""
    try:
        user_id = session.get('user_id')
        if not user_id:
            return jsonify({"error": "User session expired"}), 401
        
        data = request.json
        if not data or 'answer' not in data or 'question' not in data:
            return jsonify({"error": "Missing required parameters"}), 400
        
        user_answer = data['answer'].strip()
        question_data = data['question']
        
        if not user_answer:
            return jsonify({"error": "No answer provided"}), 400
        
        print(f"Evaluating answer for question: {question_data['question'][:50]}...")
        
        # Prepare evaluation prompt
        prompt = f"""
        You are an educational assessment AI tasked with evaluating student answers to open-ended questions.
        
        QUESTION: {question_data['question']}
        
        MODEL ANSWER: {question_data['model_answer']}
        
        EVALUATION CRITERIA: {question_data['evaluation_criteria']}
        
        STUDENT ANSWER: {user_answer}
        
        Please evaluate the student's answer on a scale of 1-10 based on the criteria provided.
        
        FORMAT YOUR RESPONSE USING MARKDOWN:
        1. Start with a numerical score (1-10).
        2. Provide a detailed evaluation explaining the strengths and weaknesses of the answer.
        3. Include specific points that were well-addressed and those that were missed.
        4. Explain how the answer could be improved to get a higher score.
        
        Your response must start with "SCORE: X" where X is the numerical score.
        Use the following format for your evaluation:
        
        SCORE: [numerical score]
        
        [Evaluation with strengths, weaknesses, and suggestions for improvement]
        """
        
        try:
            response = model.generate_content(prompt)
            evaluation_text = response.text
            
            # Extract score from response
            score_match = re.search(r'SCORE:\s*(\d+)', evaluation_text, re.IGNORECASE)
            if score_match:
                score = int(score_match.group(1))
                # Ensure score is within 1-10 range
                score = max(1, min(10, score))
            else:
                # Default score if not found
                score = 5
            
            # Remove the score line from evaluation text if present
            evaluation_text = re.sub(r'SCORE:\s*\d+\n*', '', evaluation_text, flags=re.IGNORECASE).strip()
            
            return jsonify({
                "success": True,
                "score": score,
                "evaluation": evaluation_text
            })
            
        except Exception as e:
            error_msg = f"Error evaluating answer: {str(e)}"
            print(error_msg)
            return jsonify({"error": error_msg}), 500
        
    except Exception as e:
        error_msg = f"Unexpected error in evaluate_answer: {str(e)}"
        print(error_msg)
        return jsonify({"error": error_msg}), 500

if __name__ == '__main__':
    # Only use debug mode when running directly, not with waitress
    app.run(debug=True)